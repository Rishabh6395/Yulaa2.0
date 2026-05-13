"""
Yulaa 2.0 – Feature Presentation Generator
Creates a comprehensive PowerPoint deck describing all product features in simple English.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ── Brand colours ──────────────────────────────────────────────────────────────
PRIMARY   = RGBColor(0x25, 0x5C, 0xDB)   # Yulaa blue
ACCENT    = RGBColor(0xFF, 0x6B, 0x35)   # Orange accent
DARK      = RGBColor(0x0F, 0x17, 0x2A)   # Near-black
LIGHT     = RGBColor(0xF0, 0xF4, 0xFF)   # Light blue bg
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
MUTED     = RGBColor(0x64, 0x74, 0x8B)
GREEN     = RGBColor(0x10, 0xB9, 0x81)
YELLOW    = RGBColor(0xF5, 0x9E, 0x0B)

prs = Presentation()
prs.slide_width  = Inches(16)
prs.slide_height = Inches(9)

BLANK = prs.slide_layouts[6]   # completely blank

# ── Helpers ────────────────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_color, opacity=None):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Segoe UI"
    return txBox


def add_bullet_box(slide, items, left, top, width, height,
                   font_size=14, color=DARK, bullet="●"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Segoe UI"
    return txBox


def feature_slide(title, subtitle, icon_emoji, description, bullets,
                  example_label, example_text,
                  accent_color=PRIMARY):
    """Standard feature slide layout."""
    slide = prs.slides.add_slide(BLANK)

    # left panel background
    add_rect(slide, 0, 0, 5.5, 9, accent_color)

    # right panel background
    add_rect(slide, 5.5, 0, 10.5, 9, LIGHT)

    # icon
    add_text(slide, icon_emoji, 0.3, 0.3, 2, 1.2, font_size=52, color=WHITE)

    # title
    add_text(slide, title, 0.3, 1.6, 5, 0.7, font_size=26, bold=True, color=WHITE)

    # subtitle
    add_text(slide, subtitle, 0.3, 2.4, 5, 0.5, font_size=13, color=RGBColor(0xBF,0xD7,0xFF), italic=True)

    # description
    add_text(slide, description, 0.3, 3.1, 5, 2.2, font_size=13, color=WHITE)

    # ── right panel ─────────────────────────────────────────────────────────────
    add_text(slide, "What you can do", 5.8, 0.3, 9.5, 0.5,
             font_size=11, bold=True, color=MUTED)
    add_bullet_box(slide, bullets, 5.8, 0.85, 9.7, 4.2, font_size=14, color=DARK)

    # example box
    add_rect(slide, 5.8, 5.4, 9.8, 3.2, WHITE)
    # left border accent
    add_rect(slide, 5.8, 5.4, 0.08, 3.2, accent_color)
    add_text(slide, f"📌 Example – {example_label}", 6.0, 5.45, 9.4, 0.45,
             font_size=12, bold=True, color=accent_color)
    add_text(slide, example_text, 6.0, 5.95, 9.4, 2.5,
             font_size=13, color=DARK)

    return slide


def section_divider(title, subtitle, bg_color=DARK):
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, 16, 9, bg_color)
    add_rect(slide, 0, 3.9, 16, 0.06, ACCENT)
    add_text(slide, title, 1, 3.2, 14, 1.1, font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, subtitle, 1, 4.4, 14, 0.7, font_size=18, color=RGBColor(0xA0,0xB4,0xD0),
             align=PP_ALIGN.CENTER, italic=True)
    return slide


def title_slide():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, 16, 9, DARK)
    add_rect(slide, 0, 0, 6.5, 9, PRIMARY)          # left accent band
    add_rect(slide, 6.5, 4.3, 9.5, 0.06, ACCENT)    # bottom divider

    add_text(slide, "Yulaa", 0.5, 1.0, 5.5, 1.5, font_size=80, bold=True, color=WHITE)
    add_text(slide, "2.0", 0.5, 2.5, 5.5, 1.0, font_size=60, bold=True, color=ACCENT)
    add_text(slide, "Smart School Management Platform",
             0.5, 3.6, 5.7, 0.7, font_size=15, color=RGBColor(0xBF,0xD7,0xFF), italic=True)

    add_text(slide, "Complete Feature Guide",
             7.0, 1.5, 8.5, 0.9, font_size=32, bold=True, color=WHITE)
    add_text(slide,
             "Everything a modern school needs — students, teachers, fees, attendance,\n"
             "transport, online classes, exams, and more — in one unified platform.",
             7.0, 2.7, 8.5, 2.0, font_size=15, color=RGBColor(0xA0,0xB4,0xD0))

    # platform badges
    for i, (label, col) in enumerate(
        [("Web Dashboard", PRIMARY), ("Mobile App", GREEN), ("Marketing Site", YELLOW)]):
        bx = 7.0 + i * 2.9
        add_rect(slide, bx, 5.0, 2.6, 0.55, col)
        add_text(slide, label, bx, 5.0, 2.6, 0.55, font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, "Powered by Next.js 14 · PostgreSQL · React Native · Turbo Monorepo",
             7.0, 6.8, 8.5, 0.4, font_size=11, color=MUTED)
    return slide


def overview_slide():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, 16, 9, LIGHT)
    add_rect(slide, 0, 0, 16, 1.1, PRIMARY)
    add_text(slide, "Platform Overview – 25+ Modules at a Glance", 0.4, 0.2, 15, 0.7,
             font_size=24, bold=True, color=WHITE)

    modules = [
        ("🏫", "Schools",        "Multi-school SaaS"),
        ("🎓", "Students",       "Enrol & manage"),
        ("👩‍🏫", "Teachers",      "Staff directory"),
        ("👨‍👩‍👧", "Parents",     "Child tracking"),
        ("🗓", "Attendance",      "Daily & monthly"),
        ("💰", "Fees",           "Invoices & payments"),
        ("📚", "Homework",       "Assign & submit"),
        ("📢", "Announcements",  "School notices"),
        ("✈️", "Leave",          "Apply & approve"),
        ("🚌", "Transport",      "Routes & rides"),
        ("💻", "Online Classes", "Live sessions"),
        ("📖", "Courses",        "Enrol & progress"),
        ("📝", "Exams",          "Schedule & results"),
        ("📅", "Timetable",      "Class schedule"),
        ("📋", "Syllabus",       "Subject topics"),
        ("🗓", "Events",         "School calendar"),
        ("🩺", "Compliance",     "Policy checklist"),
        ("🛒", "Vendor/Shop",    "Buy supplies"),
        ("🧑‍💼", "Consultant",   "Career guidance"),
        ("📊", "Reports",        "Export & analyse"),
        ("🔔", "Queries",        "Ask admin"),
        ("📃", "Letters",        "Custom templates"),
        ("🏗️", "Masters",        "Lookup config"),
        ("🔐", "Roles & Access", "RBAC security"),
        ("📱", "Mobile App",     "iOS & Android"),
    ]

    cols = 5
    rows = 5
    cw, rh = 3.1, 1.5
    for idx, (emoji, name, desc) in enumerate(modules):
        col = idx % cols
        row = idx // cols
        bx = 0.05 + col * cw
        by = 1.2 + row * rh
        add_rect(slide, bx, by, cw - 0.1, rh - 0.1, WHITE)
        add_text(slide, emoji, bx + 0.1, by + 0.08, 0.6, 0.5, font_size=20)
        add_text(slide, name, bx + 0.7, by + 0.1, cw - 0.8, 0.4,
                 font_size=13, bold=True, color=DARK)
        add_text(slide, desc, bx + 0.7, by + 0.55, cw - 0.8, 0.4,
                 font_size=11, color=MUTED)
    return slide


# ════════════════════════════════════════════════════════════════════════════════
# BUILD THE DECK
# ════════════════════════════════════════════════════════════════════════════════

title_slide()
overview_slide()

# ── SECTION 1 – PLATFORM FOUNDATION ──────────────────────────────────────────
section_divider("Platform Foundation",
                "How Yulaa is built to support many schools securely", PRIMARY)

feature_slide(
    title="Multi-School & Multi-Tenant",
    subtitle="One platform, many schools",
    icon_emoji="🏫",
    description="Yulaa can manage many schools at once on the same platform. Each school has its own data, users, and settings — completely separated from others.",
    bullets=[
        "Each school gets its own students, teachers, fees, and settings",
        "Super Admin can view and manage all schools from one place",
        "Schools can be 'default' for platform-level operations",
        "Each school has subscription plans (basic / premium)",
        "School profile: name, logo, address, board type, facilities",
        "Toggle features per school: online classes, courses, external vendors",
    ],
    example_label="School Setup",
    example_text=(
        "Delhi Public School – Sector 45 is added to Yulaa.\n"
        "It has its own 800 students, 50 teachers, and fee structures.\n"
        "St. Xavier's School is also on the same platform but their data is\n"
        "completely separate — admins of DPS cannot see Xavier's data."
    ),
    accent_color=PRIMARY,
)

feature_slide(
    title="Role-Based Access Control (RBAC)",
    subtitle="The right access for the right person",
    icon_emoji="🔐",
    description="Every user in Yulaa has a role. The role decides what they can see and do. This keeps data safe and operations clean.",
    bullets=[
        "Roles: Super Admin, School Admin, Teacher, Parent, Student, Vendor, Consultant",
        "A user can have different roles in different schools",
        "Menu permissions can be customised per school per role",
        "Sensitive operations (fee collection, admission approval) restricted to admins",
        "JWT-based authentication — secure, stateless login",
        "OTP-based login option for mobile users",
    ],
    example_label="Role in Action",
    example_text=(
        "Mrs. Sharma logs in as a Teacher at DPS.\n"
        "She can mark attendance, assign homework, and see her class students.\n"
        "She CANNOT access fee invoices or admission applications — those are\n"
        "visible only to School Admin and Super Admin."
    ),
    accent_color=RGBColor(0x6D,0x28,0xD9),
)

# ── SECTION 2 – PEOPLE MANAGEMENT ────────────────────────────────────────────
section_divider("People Management",
                "Students · Teachers · Parents – all in one place", RGBColor(0x06,0x4E,0x3B))

feature_slide(
    title="Student Management",
    subtitle="Complete student lifecycle from admission to alumni",
    icon_emoji="🎓",
    description="Manage every student's complete profile: personal details, class, admission number, photo, blood group, Aadhaar, and parent links — all in one place.",
    bullets=[
        "Add, edit, and search students by name, class, or status",
        "Student profile: DOB, gender, blood group, address, photo",
        "Assign students to class and section automatically",
        "Bulk import students via CSV upload",
        "Export student list to Excel / CSV",
        "Track admission number, joining date, and status (active/inactive/alumni)",
        "Link multiple students to one parent across different schools",
    ],
    example_label="Adding a New Student",
    example_text=(
        "Admin adds Arjun Mehta (Age 10, Class 5-A) to DPS.\n"
        "System assigns admission number: DPS-2024-0451.\n"
        "Parent Rakesh Mehta is linked to Arjun — Rakesh can now see\n"
        "Arjun's attendance, fees, and homework from the parent app."
    ),
    accent_color=GREEN,
)

feature_slide(
    title="Teacher Management",
    subtitle="Staff directory with qualifications and assignments",
    icon_emoji="👩‍🏫",
    description="Maintain a complete directory of all teaching and non-teaching staff. Track qualifications, experience, department, and subject assignments.",
    bullets=[
        "Add teacher profiles: name, employee ID, designation, department",
        "Record qualifications, experience years, and joining date",
        "Assign teachers to classes and subjects in the timetable",
        "Teachers can mark attendance and assign homework",
        "Bulk onboard teachers via CSV",
        "Teachers apply for leave through the system",
        "View teacher-wise timetable and workload",
    ],
    example_label="Teacher Onboarding",
    example_text=(
        "Principal adds Mr. Rajesh Kumar – Maths teacher, B.Ed + M.Sc, 8 years experience.\n"
        "He is assigned to Class 9A, 9B, and 10A for Mathematics.\n"
        "He can now mark attendance for his classes and assign Maths homework\n"
        "directly from the dashboard or mobile app."
    ),
    accent_color=RGBColor(0x05,0x96,0x68),
)

feature_slide(
    title="Parent Management",
    subtitle="Keep parents informed and connected",
    icon_emoji="👨‍👩‍👧",
    description="Parents have their own login to track their children's progress across schools. One parent account can monitor multiple children.",
    bullets=[
        "Parent portal: view attendance, fees, homework, and announcements",
        "One parent can link children from different schools",
        "Receive fee payment reminders and due date alerts",
        "Raise queries and get replies from school admin",
        "View child's leave status and academic calendar",
        "Mobile app with real-time notifications",
        "Bulk parent import via CSV",
    ],
    example_label="Parent Portal",
    example_text=(
        "Priya Verma has two children — one in DPS (Class 6) and one in Ryan International (Class 3).\n"
        "She logs into Yulaa once and sees BOTH children's dashboards.\n"
        "She pays DPS fees online and sees Ryan's attendance calendar — all from one screen."
    ),
    accent_color=RGBColor(0xDB,0x27,0x77),
)

# ── SECTION 3 – ACADEMIC OPERATIONS ──────────────────────────────────────────
section_divider("Academic Operations",
                "Attendance · Homework · Timetable · Syllabus · Exams", RGBColor(0x1E,0x3A,0x5F))

feature_slide(
    title="Attendance Management",
    subtitle="Daily, monthly, and calendar views",
    icon_emoji="🗓",
    description="Teachers mark attendance class-by-class every day. Admins get monthly reports and calendar views. Parents see their child's attendance instantly.",
    bullets=[
        "Mark Present / Absent / Late with a single tap",
        '"Mark All Present" button for quick bulk marking',
        "Edit attendance even after it is submitted",
        "Monthly attendance calendar view per student",
        "Attendance percentage calculated automatically",
        "Check-in / check-out time tracking for staff",
        "Holiday calendar integration — no attendance on holidays",
        "Export monthly attendance reports",
    ],
    example_label="Daily Attendance",
    example_text=(
        "Mrs. Sharma opens Class 9A at 9:00 AM on Monday.\n"
        "She clicks 'Mark All Present', then marks 2 students Late and 1 Absent.\n"
        "Parents of absent students get a notification immediately.\n"
        "The admin dashboard shows 94% attendance for Class 9A today."
    ),
    accent_color=YELLOW,
)

feature_slide(
    title="Homework Management",
    subtitle="Assign, track, and submit homework digitally",
    icon_emoji="📚",
    description="Teachers assign homework to a class for a subject with a due date. Students and parents see it on the app. Submission counts are tracked automatically.",
    bullets=[
        "Assign homework per class, section, and subject",
        "Set due date — overdue tasks are highlighted in red",
        "Track number of submissions received",
        "Students view pending and completed homework on the app",
        "Teachers can attach notes or instructions",
        "Parents see their child's homework from the parent portal",
        "Filter homework by class, subject, or status",
    ],
    example_label="Assigning Homework",
    example_text=(
        "Mr. Rajesh assigns: 'Solve exercises 5.1 – 5.5 from NCERT Math textbook'\n"
        "to Class 9A, due Friday.\n"
        "38 students see it on their app. By Thursday, 30 have submitted.\n"
        "Mr. Rajesh sees '30/38 submitted' and follows up with the remaining 8."
    ),
    accent_color=RGBColor(0xF5,0x9E,0x0B),
)

feature_slide(
    title="Timetable Management",
    subtitle="Weekly class schedule with teacher assignments",
    icon_emoji="📅",
    description="Create and manage the weekly timetable for every class. Assign teachers to periods. Teachers can see their personal schedule, and changes are logged.",
    bullets=[
        "Create timetable for each class: periods, subjects, teachers",
        "Drag-and-drop style period assignment",
        "View teacher-wise timetable across all classes",
        "Reassign periods when a teacher is on leave",
        "Every timetable change is logged with reason and actor",
        "Students and teachers see today's schedule on the mobile app",
        "Holiday-aware: no classes on holidays",
    ],
    example_label="Timetable in Action",
    example_text=(
        "Class 9A has: Mon-Period 1 = Maths (Mr. Rajesh), Period 2 = English (Mrs. Priya)...\n"
        "Mr. Rajesh is on leave on Wednesday. Admin reassigns his periods to Mr. Suresh.\n"
        "The change is logged: 'Reason: Medical leave. Changed by: Principal Gupta.'\n"
        "Students and Mr. Suresh both see the updated schedule on their phones."
    ),
    accent_color=RGBColor(0x06,0x78,0xB2),
)

feature_slide(
    title="Syllabus Management",
    subtitle="Topic-wise curriculum tracker",
    icon_emoji="📋",
    description="Teachers plan and track what topics have been taught. Admin and parents can see the syllabus coverage. Stay on track throughout the academic year.",
    bullets=[
        "Add syllabus items: subject, topic, chapter, target date",
        "Teachers mark topics as completed with actual completion date",
        "View syllabus coverage percentage per subject",
        "Assign syllabus items to specific teachers",
        "Colour-coded: on-track, behind schedule, completed",
        "Academic year cycle management",
        "Integrates with exam planning",
    ],
    example_label="Syllabus Tracking",
    example_text=(
        "For Class 10 Maths, the syllabus has 12 chapters planned for the year.\n"
        "By October, Mr. Rajesh has marked 7 chapters as 'completed'.\n"
        "Admin sees 58% coverage — on track for the midterm exam in November.\n"
        "Chapter 8 is flagged 'behind schedule' — admin can follow up."
    ),
    accent_color=RGBColor(0x7C,0x3A,0xED),
)

feature_slide(
    title="Exam Management",
    subtitle="Schedule exams and record results",
    icon_emoji="📝",
    description="Plan exams for classes and subjects. Record results and generate performance reports. Parents see their child's grades on the app.",
    bullets=[
        "Create exams: name, type (unit test / midterm / annual), date, class, subject",
        "Multiple exam types configurable (MCQ, written, practical, etc.)",
        "Enter marks for each student after the exam",
        "Auto-calculate grades using grading type config",
        "Performance reports per student, class, and subject",
        "Compare performance across terms",
        "Export results to PDF or Excel",
    ],
    example_label="Exam Flow",
    example_text=(
        "Admin schedules 'Unit Test 1 – Mathematics – Class 9A' for 15 Nov.\n"
        "After the exam, Mr. Rajesh enters marks: Arjun = 45/50, Priya = 38/50...\n"
        "System auto-grades: A+, B+ etc. based on the school's grading scale.\n"
        "Parents see their child's result on the mobile app that evening."
    ),
    accent_color=RGBColor(0xDC,0x26,0x26),
)

# ── SECTION 4 – FINANCE ───────────────────────────────────────────────────────
section_divider("Fee Management",
                "Invoices · Payments · Reminders · PDF Export", RGBColor(0x14,0x53,0x2D))

feature_slide(
    title="Fee Management",
    subtitle="Invoices, payments, and collection tracking",
    icon_emoji="💰",
    description="Create fee structures, generate invoices for students, record payments, and track outstanding dues — all with PDF export and payment reminders.",
    bullets=[
        "Define fee structures: tuition, transport, activity, etc.",
        "Generate individual or bulk invoices for students",
        "Record payments: cash, UPI, bank transfer — with transaction reference",
        "Status tracking: Paid / Unpaid / Partial / Overdue",
        "Summary dashboard: total collected, pending, overdue count",
        "Send fee reminders via notifications",
        "Export invoices as PDF (downloadable)",
        "Upload bulk fee data via CSV",
        "Instalment-based fee collection supported",
    ],
    example_label="Fee Collection",
    example_text=(
        "Admin generates Term 2 fees for all 800 students of DPS.\n"
        "Invoice No. INV-2024-0923 for Arjun Mehta: ₹24,500 due by 30 Nov.\n"
        "Parent pays ₹24,500 via UPI — admin records payment. Status → Paid.\n"
        "Admin can export a PDF receipt and send it to the parent."
    ),
    accent_color=GREEN,
)

# ── SECTION 5 – ADMISSION ─────────────────────────────────────────────────────
section_divider("Admissions",
                "Online Applications · Workflow Approvals · Auto-Provisioning", RGBColor(0x3B,0x1F,0x6B))

feature_slide(
    title="Admission Management",
    subtitle="From online application to student creation",
    icon_emoji="🏛️",
    description="Parents apply online for their child's admission. The school reviews the application through a configurable multi-step workflow. On approval, a student account is automatically created.",
    bullets=[
        "Public admission form — no login required for parents to apply",
        "OTP verification for parent identity during application",
        "Submit details for multiple children in one application",
        "Configurable approval workflow: e.g., Coordinator → Vice Principal → Principal",
        "Each step notifies the approver; action is logged with comments",
        "Risk scoring and validation flags on applications",
        "On approval: student record + parent account auto-created",
        "Admin can view all applications: pending, approved, rejected",
        "Search and filter applications by name, class, status",
    ],
    example_label="Admission Process",
    example_text=(
        "Parent Rakesh Mehta fills the online form for his son Arjun (applying for Grade 5).\n"
        "Step 1: Coordinator reviews → approves. Step 2: Principal reviews → approves.\n"
        "System auto-creates: Student record (Arjun, Grade 5A) + Parent login for Rakesh.\n"
        "Rakesh gets a welcome SMS with his login credentials."
    ),
    accent_color=RGBColor(0x6D,0x28,0xD9),
)

# ── SECTION 6 – LEAVE ─────────────────────────────────────────────────────────
section_divider("Leave Management",
                "Apply · Approve · Balance Tracking", RGBColor(0x1C,0x3D,0x5E))

feature_slide(
    title="Leave Management",
    subtitle="Multi-step leave approval with balance tracking",
    icon_emoji="✈️",
    description="Teachers and staff apply for leave online. The request goes through a multi-step approval workflow. Leave balances are tracked automatically per leave type.",
    bullets=[
        "Teachers apply for leave: type, date range, reason",
        "Leave types: casual, sick, earned, maternity, etc. (configurable)",
        "Multi-step approval workflow (e.g., HOD → Principal)",
        "Each approver gets a notification; can approve or reject with comment",
        "Leave balance policy: each type has a quota per year",
        "Auto-deduct from balance on approval; restore on rejection",
        "Holiday-aware: only working days counted in leave duration",
        "Admin can view all leave requests with full action history",
        "Students can also submit leave applications via parent",
    ],
    example_label="Leave Request Flow",
    example_text=(
        "Mr. Rajesh applies for Sick Leave from 10–12 Jan (3 days).\n"
        "HOD Mrs. Patel approves. Principal Mr. Sharma approves.\n"
        "Leave balance updated: Sick Leave remaining = 7 days (was 10).\n"
        "Admin sees the full log: 'Approved by HOD at 10:15 AM, Principal at 11:30 AM.'"
    ),
    accent_color=RGBColor(0x02,0x78,0xBD),
)

# ── SECTION 7 – COMMUNICATION ─────────────────────────────────────────────────
section_divider("Communication",
                "Announcements · Queries · Notifications", RGBColor(0x1A,0x20,0x2C))

feature_slide(
    title="Announcements",
    subtitle="School-wide and targeted communications",
    icon_emoji="📢",
    description="School admin posts announcements for different audiences — all students, specific classes, teachers only, or parents. Important notices are highlighted as urgent.",
    bullets=[
        "Create announcements with title, message, and type",
        "Types: General, Urgent, Event, Holiday, Exam, Fee Reminder",
        "Target specific audience: all, teachers, parents, students, specific class",
        "Students, parents, and teachers see announcements on their dashboard/app",
        "Urgent announcements shown with a red banner",
        "Filter and search past announcements",
        "Customisable announcement types via Masters",
    ],
    example_label="Posting an Announcement",
    example_text=(
        "Admin posts: 'School closed on 26 Jan for Republic Day' — Type: Holiday.\n"
        "All 800 students' parents and 50 teachers receive a push notification.\n"
        "Admin also posts: 'Term 2 exams from 15 Nov — Type: Exam' targeted\n"
        "only to Class 9 and 10 students and their parents."
    ),
    accent_color=ACCENT,
)

feature_slide(
    title="Queries & Support",
    subtitle="Two-way communication between parents and school",
    icon_emoji="🔔",
    description="Parents or students raise queries to the school administration. Admin responds and the thread is tracked. Priority levels ensure urgent issues are handled first.",
    bullets=[
        "Parents raise queries: fee doubt, admission status, class issue, etc.",
        "Admin sees all queries with priority levels (low/medium/high/urgent)",
        "Admin responds directly — parent sees the reply on their app",
        "Full conversation thread maintained per query",
        "Status tracking: open, in-progress, resolved",
        "Support queries (admin-to-platform) tracked separately",
        "Filter queries by status, priority, or date",
    ],
    example_label="Query in Action",
    example_text=(
        "Parent Rakesh Mehta raises a query: 'Why was Arjun marked absent on 5th Nov? He was present.'\n"
        "Priority: Medium. Admin checks the attendance record and responds:\n"
        "'Apologies, error corrected — attendance updated to Present.'\n"
        "Rakesh sees the reply and closes the query."
    ),
    accent_color=RGBColor(0x0E,0x7A,0x0C),
)

# ── SECTION 8 – TRANSPORT ─────────────────────────────────────────────────────
section_divider("Transport Management",
                "Routes · Buses · Daily Rides · Student Subscriptions", RGBColor(0x27,0x1D,0x0A))

feature_slide(
    title="Transport Management",
    subtitle="School bus routes, buses, and daily rides",
    icon_emoji="🚌",
    description="Manage school buses, define routes, assign drivers, and track daily rides. Students subscribe to a route and their pick-up/drop status is tracked in real time.",
    bullets=[
        "Define transport routes with pickup points and estimated times",
        "Register school buses: registration number, capacity, driver",
        "Students subscribe to a route (morning pickup / evening drop)",
        "Record daily rides: start and end times, passengers",
        "Driver / route operator assigned per ride",
        "Admin dashboard shows today's active rides and route status",
        "Transport reports: monthly ridership, route utilisation",
        "Integrate transport fees with the fee management module",
    ],
    example_label="Bus Route Setup",
    example_text=(
        "Route 'Sector 12 → DPS' covers 8 stops, 45 students subscribed.\n"
        "Bus: DL-1P-4532, Driver: Ramesh Kumar.\n"
        "Monday morning: Ramesh marks ride started at 7:00 AM.\n"
        "Admin sees: 45 students on board, ETA 8:10 AM. Ride completed at 8:12 AM."
    ),
    accent_color=RGBColor(0xD9,0x77,0x07),
)

# ── SECTION 9 – ONLINE CLASSES & COURSES ──────────────────────────────────────
section_divider("Online Classes & Courses",
                "Live Sessions · Recorded Courses · Enrollments & Progress", RGBColor(0x0C,0x22,0x3A))

feature_slide(
    title="Online Classes",
    subtitle="Schedule and conduct live virtual sessions",
    icon_emoji="💻",
    description="Teachers schedule live online classes for their students. Students join via the app or web. Attendance is tracked automatically for each session.",
    bullets=[
        "Create an online class: title, subject, class, date, time, duration",
        "Integration with external platforms (Zoom, Google Meet, etc.)",
        "Students see upcoming classes on their dashboard/app",
        "Live attendance tracking: who joined and when",
        "Attendance report per session: present/absent per student",
        "Admin configures which external platform to use",
        "Super Admin can enable/disable online classes per school",
        "Teachers manage their own online class schedule",
    ],
    example_label="Live Class Session",
    example_text=(
        "Mr. Rajesh schedules: 'Online Maths Revision – Class 10A'\n"
        "for Saturday 10 AM, 60 minutes, via Google Meet.\n"
        "38 students see the class in their app with a 'Join' button.\n"
        "32 students join. Attendance auto-marked — 6 students absent from session."
    ),
    accent_color=PRIMARY,
)

feature_slide(
    title="Courses",
    subtitle="Self-paced learning with modules and progress tracking",
    icon_emoji="📖",
    description="Schools and the platform offer courses students can enrol in. Courses are structured into modules. Progress is tracked as students complete each module.",
    bullets=[
        "Create courses: title, description, duration, teacher, category",
        "Courses have multiple modules / lessons in sequence",
        "Students enrol in a course (free or fee-based)",
        "Progress tracker: percentage completed per student",
        "Super Admin manages platform-wide courses",
        "School admin enables courses for their school",
        "Students see enrolled courses and continue from where they left off",
        "Teachers manage the courses they own",
    ],
    example_label="Course Enrollment",
    example_text=(
        "Platform offers: 'JEE Preparation – Physics' — 12 modules, 40 hours.\n"
        "Arjun (Class 11) enrolls. He completes Modules 1–4 over two weeks.\n"
        "Dashboard shows: 33% complete. Module 5 is next: 'Laws of Motion.'\n"
        "Teacher sees 45 students enrolled, 12 have completed > 50%."
    ),
    accent_color=RGBColor(0x4F,0x46,0xE5),
)

# ── SECTION 10 – EVENTS & COMPLIANCE ─────────────────────────────────────────
section_divider("Events, Compliance & More",
                "School Calendar · Compliance · Vendor Shop · Career Sessions", RGBColor(0x1A,0x1A,0x2E))

feature_slide(
    title="School Events",
    subtitle="Plan and manage school activities and celebrations",
    icon_emoji="🎉",
    description="Create school events like sports day, annual day, parent-teacher meetings, and field trips. Assign tasks to staff and track event preparation.",
    bullets=[
        "Create events: name, type, date, description, venue",
        "Customisable event types (Annual Day, Sports, PTM, Excursion, etc.)",
        "Assign tasks to teachers / staff for event preparation",
        "Events appear on the school calendar for all users",
        "Parents see upcoming events on the mobile app",
        "Announcements can be linked to events",
        "Event history for the academic year",
    ],
    example_label="Annual Day Planning",
    example_text=(
        "Admin creates event: 'Annual Day 2025 – 15 March 2025, School Auditorium'.\n"
        "Tasks assigned: Mrs. Priya → cultural programme, Mr. Ravi → sound system.\n"
        "All parents see the event in their app calendar and get a reminder 3 days before.\n"
        "Post-event: Admin archives it with photos and notes."
    ),
    accent_color=ACCENT,
)

feature_slide(
    title="Compliance Management",
    subtitle="Track school policy and regulatory requirements",
    icon_emoji="🩺",
    description="Schools must meet various compliance requirements — safety drills, document submissions, inspections. Yulaa helps track status of each compliance item.",
    bullets=[
        "Create compliance checklist items with due dates",
        "Track status: pending / in-progress / completed / overdue",
        "Assign responsibility to a specific staff member",
        "Admin gets an overview of all compliance items and their status",
        "Colour-coded dashboard: green (done), red (overdue), yellow (in-progress)",
        "Add notes and attach evidence (document path / URL)",
        "Per-school compliance management",
    ],
    example_label="Compliance Tracking",
    example_text=(
        "School has 10 compliance items for the year:\n"
        "✅ Fire drill (Apr) — Completed | ⚠️ CBSE affiliation renewal (Sep) — Due in 15 days\n"
        "❌ Building safety certificate — Overdue by 10 days\n"
        "Admin sees a red alert and assigns Mr. Gupta to resolve it immediately."
    ),
    accent_color=RGBColor(0x06,0x69,0x2C),
)

feature_slide(
    title="Vendor / School Shop",
    subtitle="Buy school supplies and uniforms online",
    icon_emoji="🛒",
    description="Vendors (uniform suppliers, book sellers, stationery shops) list their products on Yulaa. Parents browse and order from the school-approved shop. Vendors manage their inventory and orders.",
    bullets=[
        "Vendors register and list products with price and stock",
        "Parents browse and place orders for school supplies",
        "Vendors see incoming orders and mark them fulfilled",
        "Admin approves which vendors can operate in their school",
        "Parents can rate and review vendors",
        "Vendor inventory management: add, edit, deactivate products",
        "Order history for parents: past purchases",
        "Super Admin manages platform-wide vendor approvals",
    ],
    example_label="Vendor Shop",
    example_text=(
        "Sharma Uniforms lists DPS school uniform set at ₹850.\n"
        "Parent Rakesh orders 2 sets for Arjun's new school year.\n"
        "Sharma Uniforms sees the order, packs it, marks 'Dispatched'.\n"
        "Rakesh receives the order and gives a 5-star rating."
    ),
    accent_color=RGBColor(0xE5,0x7B,0x08),
)

feature_slide(
    title="Career Sessions (Consultant)",
    subtitle="Book one-on-one sessions with career counsellors",
    icon_emoji="🧑‍💼",
    description="External career consultants offer guidance sessions to students and parents. They set their availability, and users can book slots directly through Yulaa.",
    bullets=[
        "Consultants register and set their weekly availability",
        "Students/parents browse consultants and book available slots",
        "Session confirmation sent to both parties",
        "School can allow only specific consultants (contract-based)",
        "Consultants manage their upcoming and past sessions",
        "Users rate sessions after completion",
        "Super Admin oversees all consultants on the platform",
        "Consultant contracts define which schools they serve",
    ],
    example_label="Booking a Career Session",
    example_text=(
        "Mr. Anand Kapoor (IIT Counsellor) has slots available: Mon/Wed/Fri 4–6 PM.\n"
        "Arjun (Class 11) books a 30-min slot on Monday 4:30 PM.\n"
        "Both receive a confirmation. After the session, Arjun rates it 5 stars.\n"
        "Admin sees total sessions booked this month: 42."
    ),
    accent_color=RGBColor(0x7C,0x3A,0xED),
)

# ── SECTION 11 – REPORTS & ADMIN ─────────────────────────────────────────────
section_divider("Reports, Letters & Admin Tools",
                "Data Export · Letter Templates · Master Configuration", RGBColor(0x0D,0x1B,0x2A))

feature_slide(
    title="Reports & Analytics",
    subtitle="Monthly reports and data exports",
    icon_emoji="📊",
    description="Generate detailed reports on attendance, fees, students, and performance. Export them as Excel or PDF files for record-keeping or inspection.",
    bullets=[
        "Monthly attendance report: class-wise, student-wise",
        "Fee collection report: collected vs pending vs overdue",
        "Student performance report: exam-wise marks and grades",
        "Export any report to Excel (.xlsx) or PDF",
        "Dashboard summary: real-time stats on one screen",
        "Scheduled/cron-based data refresh for the dashboard",
        "Redis caching for fast dashboard load",
        "Admin can drill down: school → class → student",
    ],
    example_label="Monthly Report",
    example_text=(
        "At the end of October, Principal downloads the Monthly Attendance Report for Class 9A.\n"
        "Report shows: 35/38 students had > 85% attendance. 3 students flagged with < 75%.\n"
        "He also exports Fee Collection Report: ₹18,50,000 collected, ₹2,20,000 pending.\n"
        "Both reports are saved as Excel files and shared with the trust board."
    ),
    accent_color=PRIMARY,
)

feature_slide(
    title="Letter Templates",
    subtitle="Create custom official letters",
    icon_emoji="📃",
    description="Schools can create standard letter templates for TC (Transfer Certificate), bonafide, fee receipts, and other official documents. Fill in details and print.",
    bullets=[
        "Create reusable letter templates with placeholders",
        "Templates for: Transfer Certificate, Bonafide, Experience Letter, etc.",
        "Fill in student/teacher details to generate a letter",
        "Download or print the generated letter as PDF",
        "Admin manages all templates from one screen",
        "Templates can be customised per school branding",
    ],
    example_label="Transfer Certificate",
    example_text=(
        "Arjun is moving to another city. Admin opens the TC template.\n"
        "System auto-fills: student name, admission no., class, date of joining, date of leaving.\n"
        "Admin reviews, adds the principal's signature stamp, and downloads the PDF.\n"
        "TC is ready in under 2 minutes."
    ),
    accent_color=MUTED,
)

feature_slide(
    title="Masters / Lookup Configuration",
    subtitle="Customise dropdown values and school-specific lists",
    icon_emoji="🏗️",
    description="Every dropdown in the system — blood groups, genders, exam types, leave types, grades, streams — can be customised by the admin. This avoids hard-coded lists.",
    bullets=[
        "Configurable masters: Gender, Blood Group, Qualifications, Streams",
        "Grades & Grading Types (A+, A, B+ … or percentage-based)",
        "Leave Types (Casual, Sick, Earned, Maternity, etc.)",
        "Exam Types (Unit Test, Midterm, Annual, Practical)",
        "Event Types, Announcement Types, Content Types",
        "Location masters: Countries, States, Districts",
        "School Hierarchy & Location masters",
        "Custom / Generic masters for unique school needs",
    ],
    example_label="Custom Masters",
    example_text=(
        "DPS uses CBSE grading: A1, A2, B1, B2, C1, C2, D, E.\n"
        "St. Xavier's uses percentage-based grading: 90%+ = A+, 80–89% = A...\n"
        "Each school sets up its OWN grading type in Masters.\n"
        "When exam results are entered, the correct grading formula applies automatically."
    ),
    accent_color=RGBColor(0x57,0x53,0x4E),
)

# ── SECTION 12 – MOBILE APP ───────────────────────────────────────────────────
section_divider("Mobile App",
                "iOS & Android – Built with React Native / Expo", RGBColor(0x06,0x2A,0x3F))

feature_slide(
    title="Mobile App",
    subtitle="All key features on iOS and Android",
    icon_emoji="📱",
    description="Yulaa has a full-featured mobile app for teachers, parents, students, and vendors. Built with React Native (Expo) — one codebase, two platforms.",
    bullets=[
        "Teachers: mark attendance, assign homework, view timetable",
        "Parents: see child's attendance, fees due, announcements",
        "Students: view homework, online classes, courses, events",
        "Vendors: manage products and view orders",
        "OTP-based login for mobile users",
        "Push notifications for fees due, attendance, announcements",
        "Leave application from mobile",
        "Career session booking on the go",
        "Offline-friendly design for slow connectivity",
    ],
    example_label="Parent Morning Routine",
    example_text=(
        "7:45 AM: Parent Rakesh checks the app before dropping Arjun at school.\n"
        "He sees: 'Today's homework: Maths Ex 5.3 due tomorrow' + 'Fees due: ₹24,500 by 30 Nov'.\n"
        "He also sees an announcement: 'PTM on Saturday 9 AM'.\n"
        "He taps 'Pay Fees' and settles the dues in 30 seconds — done before reaching school."
    ),
    accent_color=RGBColor(0x0E,0xA5,0xE9),
)

# ── SECTION 13 – SUPER ADMIN ──────────────────────────────────────────────────
section_divider("Super Admin Panel",
                "Platform-level control over all schools and users", RGBColor(0x1A,0x0A,0x2E))

feature_slide(
    title="Super Admin Dashboard",
    subtitle="God-mode control over the entire platform",
    icon_emoji="🛡️",
    description="The Super Admin manages all schools, users, consultants, vendors, and platform-level settings from a single unified dashboard.",
    bullets=[
        "Add, edit, and deactivate schools",
        "Manage all users across all schools",
        "Approve and manage platform consultants",
        "Approve and manage platform vendors",
        "Configure online class providers (Zoom/Meet credentials per school)",
        "Manage platform-wide courses",
        "View all schools' subscription plans",
        "Reset any user's password",
        "Access all school settings and module toggles",
    ],
    example_label="Platform Management",
    example_text=(
        "Super Admin Aarav logs in and sees 24 schools, 12,000 students, 800 teachers on the platform.\n"
        "He approves a new vendor: 'Kumar Books' to supply books to 5 schools.\n"
        "He enables the Courses module for DPS (currently disabled).\n"
        "He resets the password for a school admin who is locked out — done in 30 seconds."
    ),
    accent_color=RGBColor(0x7C,0x3A,0xED),
)

# ── CLOSING SLIDE ─────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 16, 9, DARK)
add_rect(slide, 0, 0, 16, 0.08, PRIMARY)
add_rect(slide, 0, 8.92, 16, 0.08, ACCENT)

add_text(slide, "Yulaa 2.0", 1, 1.5, 14, 1.4, font_size=64, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "One Platform. Every School Need.",
         1, 3.1, 14, 0.8, font_size=26, color=RGBColor(0xA0,0xB4,0xD0),
         align=PP_ALIGN.CENTER, italic=True)

summary = [
    "25+ modules covering every aspect of school management",
    "Multi-tenant SaaS — scale from 1 school to 1,000 schools",
    "Web dashboard + Mobile App (iOS & Android)",
    "Role-based access: Admin · Teacher · Parent · Student · Vendor · Consultant",
    "Built on Next.js 14, PostgreSQL, React Native, Turborepo",
]
add_bullet_box(slide, summary, 3, 4.2, 10, 3.2,
               font_size=16, color=RGBColor(0xC8,0xD8,0xF0), bullet="✦")

add_text(slide, "© 2025 Sequel One Solutions Pvt. Ltd. — Confidential",
         1, 8.4, 14, 0.4, font_size=11, color=MUTED, align=PP_ALIGN.CENTER)

# ── SAVE ─────────────────────────────────────────────────────────────────────
OUTPUT = r"D:\OneDrive - Sequel One Solutions Pvt. Ltd\Desktop\tacker\Yulaa2.0\Yulaa2.0_Features.pptx"
prs.save(OUTPUT)
print(f"✅  Saved: {OUTPUT}")
print(f"    Slides: {len(prs.slides)}")
